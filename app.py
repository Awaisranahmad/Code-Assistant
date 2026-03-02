import streamlit as st
from groq import Groq
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
import plotly.graph_objects as go
import numpy as np
import re
import textstat
from pptx import Presentation
import io

# --- 1. Connection ---
client = Groq(api_key=st.secrets["GROQ_API_KEY"])

st.set_page_config(page_title="AI Sentinel Ultra", page_icon="🛡️", layout="wide")

# --- 2. Professional Blue & White UI ---
st.markdown("""
<style>
    .stApp { background-color: #f8fafc; color: #1e3a8a; }
    .title { font-size: 38px; font-weight: 800; text-align: center; color: #1e3a8a; margin-bottom: 20px; }
    .metric-card { background: white; padding: 20px; border-radius: 12px; box-shadow: 0 4px 12px rgba(0,0,0,0.05); text-align: center; border-top: 4px solid #3b82f6; }
    .report-card { background: white; padding: 25px; border-radius: 15px; border-left: 10px solid #1e3a8a; box-shadow: 0 10px 25px rgba(0,0,0,0.05); color: #334155; }
    .stTextArea textarea { border-radius: 12px; border: 2px solid #cbd5e1 !important; }
    [data-testid="stFileUploadDropzone"] { background: white; border: 2px dashed #3b82f6 !important; border-radius: 12px; }
</style>
""", unsafe_allow_html=True)

# --- 3. Universal File Reader (Supports EVERYTHING) ---
def read_any_file(file):
    try:
        fname = file.name.lower()
        if fname.endswith('.pdf'):
            return " ".join([p.extract_text() for p in PdfReader(file).pages])
        elif fname.endswith('.docx'):
            return " ".join([p.text for p in Document(file).paragraphs])
        elif fname.endswith('.pptx'):
            prs = Presentation(file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return " ".join(text_runs)
        elif fname.endswith(('.csv', '.xlsx')):
            df = pd.read_csv(file) if fname.endswith('.csv') else pd.read_excel(file)
            return df.to_string()
        else:
            # Code files (.py, .js, .cpp), Text files (.txt, .log, .md)
            return file.read().decode("utf-8", errors="ignore")
    except Exception as e:
        return f"System Error reading file: {str(e)}"

# --- 4. Linguistic Style Metrics ---
def style_metrics(text):
    sentences = re.split(r'[.!?]', text)
    words = text.split()
    sentence_lengths = [len(s.split()) for s in sentences if len(s.split()) > 0]
    
    avg_sentence = np.mean(sentence_lengths) if sentence_lengths else 0
    variance = np.var(sentence_lengths) if sentence_lengths else 0
    diversity = len(set(words)) / len(words) if words else 0
    readability = textstat.flesch_reading_ease(text)
    
    return {"avg_sentence": avg_sentence, "variance": variance, "diversity": diversity, "readability": readability}

# --- 5. Main UI Flow ---
st.markdown("<div class='title'>🛡️ AI Sentinel Ultra: Universal Auditor</div>", unsafe_allow_html=True)

# Step 1: Input Hierarchy
st.markdown("### ⌨️ Step 1: Paste Content (Optional)")
manual_input = st.text_area("Paste code, logs, or text", height=200, label_visibility="collapsed")

st.markdown("### 📂 Step 2: Or Upload Any File")
uploaded_file = st.file_uploader("Upload PDF, Word, Excel, PowerPoint, Python, JS, etc.", type=None)

content = ""
if uploaded_file:
    content = read_any_file(uploaded_file)
    st.info(f"✅ Active Source: {uploaded_file.name}")
else:
    content = manual_input

st.write("---")
run_btn = st.button("🚀 EXECUTE FULL UNIVERSAL SCAN", use_container_width=True)

# --- 6. Analysis Execution ---
if run_btn and content:
    with st.spinner("Analyzing cross-platform linguistic patterns..."):
        try:
            metrics = style_metrics(content)
            
            # AI Engine Deep Analysis
            prompt = (
                "Identify if this text is AI or Human. You MUST name the specific engine (e.g. GPT-4, Claude 3, Gemini, Llama 3). "
                "Format: AI_SCORE: [0-100], ENGINE: [Name], CONFIDENCE: [Score], REPORT: [One Paragraph]"
                f"\n\nCONTENT:\n{content[:3800]}"
            )
            
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.1
            )
            
            out = response.choices[0].message.content
            
            # Parsing with Fallbacks
            ai_val = int(re.search(r"AI_SCORE:\s*(\d+)", out).group(1)) if re.search(r"AI_SCORE:\s*(\d+)", out) else 50
            engine = re.search(r"ENGINE:\s*(.*)", out).group(1) if re.search(r"ENGINE:\s*(.*)", out) else "Undetermined"
            conf = re.search(r"CONFIDENCE:\s*(.*)", out).group(1) if re.search(r"CONFIDENCE:\s*(.*)", out) else "Medium"
            report_body = re.search(r"REPORT:\s*(.*)", out, re.S).group(1) if re.search(r"REPORT:\s*(.*)", out, re.S) else out

            # --- Layout Step 3: Visualization ---
            
            st.markdown("### 📊 Step 3: Audit Visualization")
            v1, v2 = st.columns([1, 1.2])

            with v1:
                fig = go.Figure(data=[go.Pie(
                    labels=['AI Signatures', 'Human Logic'], values=[ai_val, 100-ai_val], hole=.7,
                    marker_colors=['#1e3a8a', '#60a5fa']
                )])
                fig.update_layout(showlegend=True, height=350, legend=dict(orientation="h", x=0.2, y=-0.1))
                st.plotly_chart(fig, use_container_width=True)

            with v2:
                fig2 = go.Figure()
                fig2.add_bar(
                    x=["Variance", "Vocabulary (%)", "Readability"],
                    y=[metrics["variance"], metrics["diversity"]*100, metrics["readability"]],
                    marker_color='#3b82f6'
                )
                fig2.update_layout(title="Writing Pattern Analysis", height=350)
                st.plotly_chart(fig2, use_container_width=True)

            # --- Step 4: Final Report ---
            st.write("---")
            st.markdown("### 📝 Step 4: Forensic Report")
            
            c1, c2, c3 = st.columns(3)
            with c1: st.markdown(f"<div class='metric-card'><b>Detected Engine</b><br><h3>{engine}</h3></div>", unsafe_allow_html=True)
            with c2: st.markdown(f"<div class='metric-card'><b>Confidence</b><br><h3>{conf}</h3></div>", unsafe_allow_html=True)
            with c3: st.markdown(f"<div class='metric-card'><b>Vocabulary</b><br><h3>{round(metrics['diversity']*100, 1)}%</h3></div>", unsafe_allow_html=True)

            st.markdown(f"<div class='report-card'>{report_body}</div>", unsafe_allow_html=True)

        except Exception as e:
            st.error(f"Audit Error: {e}")
elif run_btn:
    st.warning("Please provide some data to analyze!")

st.write("---")
st.caption("AI Sentinel Ultra v9.0 | Universal Multi-Format Auditor")
